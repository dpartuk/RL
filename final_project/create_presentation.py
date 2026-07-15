"""Generate TED-style presentation for Messy-Network Trajectory Follower."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
PLOTS_DIR = os.path.join(OUTPUT_DIR, "experiment_1d", "output_plots")
IRL_DIAGRAM_PATH = os.path.join(OUTPUT_DIR, "experiment_1d", "output_plots", "irl_diagram.png")
OUT_PATH = os.path.join(OUTPUT_DIR, "presentation.pptx")

# --- Colors ---
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MEDIUM = RGBColor(0x22, 0x22, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0x8B)
ACCENT_BLUE = RGBColor(0x4D, 0x9D, 0xE8)
ACCENT_YELLOW = RGBColor(0xF5, 0xC5, 0x42)
ACCENT_ORANGE = RGBColor(0xF0, 0x8C, 0x3A)
DIM_WHITE = RGBColor(0xAA, 0xAA, 0xAA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=24,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox, tf


def add_paragraph(tf, text, font_size=24, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Calibri",
                  space_before=0, line_spacing=1.2):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.line_spacing = Pt(int(font_size * line_spacing))
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     text="", font_size=18, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Calibri"
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(2),
             "What If Your Robot\nCan't Trust Its Own Eyes?",
             font_size=52, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
             "Reinforcement Learning Under Data Dropouts",
             font_size=28, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Doron Peleg  |  Afeka College of Engineering, Tel Aviv",
             font_size=20, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.5),
             "Inspired by Fan et al. (2025) — IEEE Transactions on Cybernetics",
             font_size=16, color=DIM_WHITE,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: The Scenario — Opening Hook
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "Imagine this...",
             font_size=40, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(1.8), Inches(10), Inches(4),
             "A surgical robot is operating remotely.\n"
             "The Wi-Fi drops for half a second.\n\n"
             "The robot's last known position of the scalpel\n"
             "is 2mm from where it actually is.\n\n"
             "Does it cut?  Does it freeze?  Does it guess?",
             font_size=30, color=WHITE, alignment=PP_ALIGN.LEFT,
             line_spacing=1.5)

add_text_box(slide, Inches(1), Inches(6.0), Inches(10), Inches(0.8),
             "This is the data dropout problem in networked control systems.",
             font_size=22, color=ACCENT_RED, bold=True)

# ============================================================
# SLIDE 3: The Core Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2),
             "Every RL algorithm assumes the agent can see.",
             font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.8),
             "What happens when it can't?",
             font_size=36, color=ACCENT_RED, bold=True,
             alignment=PP_ALIGN.CENTER)

# Three reaction boxes
box_w = Inches(3.2)
box_h = Inches(3.0)
gap = Inches(0.5)
start_x = Inches(1.2)
y = Inches(3.2)

# Box 1 - Pretend
shape = add_rounded_rect(slide, start_x, y, box_w, box_h, RGBColor(0x8B, 0x30, 0x30))
tf = shape.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = "Pretend\nnothing happened"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nUse the stale observation\nas if it were current"
p2.font.size = Pt(16)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

# Box 2 - Predict
shape = add_rounded_rect(slide, start_x + box_w + gap, y, box_w, box_h,
                         RGBColor(0x8B, 0x6B, 0x20))
tf = shape.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = "Try to predict\nwhat you missed"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nExtrapolate from last\nknown state + action"
p2.font.size = Pt(16)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

# Box 3 - Admit
shape = add_rounded_rect(slide, start_x + 2 * (box_w + gap), y, box_w, box_h,
                         RGBColor(0x20, 0x6B, 0x4E))
tf = shape.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = "Admit you\ndon't know"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nSkip learning when\nthe data is unreliable"
p2.font.size = Pt(16)
p2.font.color.rgb = LIGHT_GRAY
p2.font.name = "Calibri"
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

# ============================================================
# SLIDE 4: Our Simplified World
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Our Simplified World",
             font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.LEFT)

# Left side: bullet points
_, tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                     "A 1D number line from -4 to +4",
                     font_size=22, color=WHITE)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "Target moves in a sine wave",
              font_size=22, color=WHITE, space_before=8)
add_paragraph(tf, "  2.0 * sin(2πt / 100)",
              font_size=20, color=ACCENT_BLUE, space_before=4)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "3 actions: LEFT (-0.2), STAY (0), RIGHT (+0.2)",
              font_size=22, color=WHITE, space_before=8)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "41 discrete states (bins of relative distance)",
              font_size=22, color=WHITE, space_before=8)
add_paragraph(tf, "  Bin 20 = right on target",
              font_size=20, color=ACCENT_GREEN, space_before=4)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "Reward = -|distance|",
              font_size=22, color=WHITE, space_before=8)
add_paragraph(tf, "  On target = 0 (best),  4 units away = -4 (worst)",
              font_size=20, color=ACCENT_RED, space_before=4)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "Each episode = 300 steps (3 sine cycles)",
              font_size=22, color=WHITE, space_before=8)

# Right side: number line diagram
# Simple visual representation
add_rounded_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5),
                 RGBColor(0x2A, 0x2A, 0x45))

add_text_box(slide, Inches(7.3), Inches(2.0), Inches(5), Inches(0.6),
             "The Number Line",
             font_size=20, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

# Number line labels
add_text_box(slide, Inches(7.5), Inches(2.8), Inches(0.6), Inches(0.4),
             "-4", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.5), Inches(2.8), Inches(0.8), Inches(0.4),
             "target", font_size=16, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(11.8), Inches(2.8), Inches(0.6), Inches(0.4),
             "+4", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Number line itself
slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(3.35), Inches(4.5), Inches(0.06)
).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = LIGHT_GRAY
slide.shapes[-1].line.fill.background()

# Target marker
add_text_box(slide, Inches(9.6), Inches(3.5), Inches(1), Inches(0.5),
             "▲", font_size=24, color=ACCENT_GREEN,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.3), Inches(3.9), Inches(1.5), Inches(0.4),
             "Target (moving)", font_size=14, color=ACCENT_GREEN,
             alignment=PP_ALIGN.CENTER)

# Agent marker
add_text_box(slide, Inches(8.6), Inches(3.5), Inches(1), Inches(0.5),
             "▲", font_size=24, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.3), Inches(3.9), Inches(1.5), Inches(0.4),
             "Agent", font_size=14, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)

# Reward examples
add_text_box(slide, Inches(7.3), Inches(4.6), Inches(5), Inches(0.4),
             "Reward examples:", font_size=18, color=ACCENT_YELLOW, bold=True)
add_text_box(slide, Inches(7.3), Inches(5.1), Inches(5), Inches(1.2),
             "On target:   reward =  0.0\n"
             "1 unit away:  reward = -1.0\n"
             "4 units away: reward = -4.0",
             font_size=16, color=LIGHT_GRAY, font_name="Courier New")

# ============================================================
# SLIDE 5: The Dropout Channel
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "The Dropout Channel: 30% Packet Loss",
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1),
             "Like a GPS that freezes 30% of the time—\n"
             "showing where you were instead of where you are.",
             font_size=26, color=LIGHT_GRAY)

# Timeline visualization
box_y = Inches(3.2)
step_w = Inches(1.1)
step_h = Inches(2.8)
start = Inches(0.8)

steps_data = [
    ("Step 1", "Bin 18", "Bin 18", True, "OK"),
    ("Step 2", "Bin 19", "Bin 19", True, "OK"),
    ("Step 3", "Bin 20", "Bin 19", False, "DROP"),
    ("Step 4", "Bin 21", "Bin 21", True, "OK"),
    ("Step 5", "Bin 22", "Bin 22", True, "OK"),
    ("Step 6", "Bin 21", "Bin 22", False, "DROP"),
    ("Step 7", "Bin 20", "Bin 20", True, "OK"),
    ("Step 8", "Bin 19", "Bin 19", True, "OK"),
    ("Step 9", "Bin 20", "Bin 20", True, "OK"),
    ("Step 10", "Bin 21", "Bin 20", False, "DROP"),
]

for i, (label, true_s, obs_s, ok, status) in enumerate(steps_data):
    x = start + Inches(i * 1.22)
    bg_c = RGBColor(0x20, 0x6B, 0x4E) if ok else RGBColor(0x8B, 0x30, 0x30)

    shape = add_rounded_rect(slide, x, box_y, step_w, step_h, bg_c)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_YELLOW
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"\nTrue:\n{true_s}"
    p2.font.size = Pt(13)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = f"\nSeen:\n{obs_s}"
    p3.font.size = Pt(13)
    p3.font.color.rgb = ACCENT_GREEN if ok else ACCENT_RED
    p3.font.bold = True
    p3.font.name = "Calibri"
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = f"\n{status}"
    p4.font.size = Pt(14)
    p4.font.color.rgb = WHITE if ok else ACCENT_RED
    p4.font.bold = True
    p4.font.name = "Calibri"
    p4.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.6),
             "Red steps: agent thinks it's in one bin but is actually in another. "
             "It makes decisions based on wrong information.",
             font_size=18, color=ACCENT_RED)

# ============================================================
# SLIDE 6: Three Agents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Three Agents, One Environment",
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Same Q-learning algorithm (41 states x 3 actions). Only difference: how they handle dropout.",
             font_size=20, color=LIGHT_GRAY)

# Three agent boxes
agent_w = Inches(3.5)
agent_h = Inches(4.2)
gap_x = Inches(0.4)
start_x = Inches(0.9)
agent_y = Inches(2.0)

# Baseline
shape = add_rounded_rect(slide, start_x, agent_y, agent_w, agent_h,
                         RGBColor(0x1E, 0x5C, 0x3A))
tf = shape.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = "BASELINE"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_GREEN
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
for line, sz, clr in [
    ("\nNo dropout at all", 20, WHITE),
    ("\nSees true state every step", 18, LIGHT_GRAY),
    ("\nAlways updates Q-table", 18, LIGHT_GRAY),
    ("\n\nThe upper bound—", 18, ACCENT_YELLOW),
    ("what's possible with\nperfect information", 18, ACCENT_YELLOW),
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(sz)
    p2.font.color.rgb = clr
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# Naive
shape = add_rounded_rect(slide, start_x + agent_w + gap_x, agent_y,
                         agent_w, agent_h, RGBColor(0x7A, 0x28, 0x28))
tf = shape.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = "NAIVE"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_RED
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
for line, sz, clr in [
    ("\n30% dropout", 20, WHITE),
    ("\nReceives stale states", 18, LIGHT_GRAY),
    ("\nAlways updates Q-table", 18, ACCENT_RED),
    ("\n\nWrites updates to wrong", 18, ACCENT_RED),
    ("Q-table entries—", 18, ACCENT_RED),
    ("corrupting its knowledge", 18, ACCENT_RED),
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(sz)
    p2.font.color.rgb = clr
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# Dropout-Aware
shape = add_rounded_rect(slide, start_x + 2 * (agent_w + gap_x), agent_y,
                         agent_w, agent_h, RGBColor(0x1E, 0x3A, 0x5C))
tf = shape.text_frame
tf.clear()
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.TOP
p = tf.paragraphs[0]
p.text = "DROPOUT-AWARE"
p.font.size = Pt(28)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER
for line, sz, clr in [
    ("\n30% dropout", 20, WHITE),
    ("\nReceives stale states", 18, LIGHT_GRAY),
    ("\nSkips Q-update if stale", 18, ACCENT_BLUE),
    ("\n\nProtects Q-table from", 18, ACCENT_BLUE),
    ("corruption—learns less", 18, ACCENT_BLUE),
    ("often, but learns correctly", 18, ACCENT_BLUE),
]:
    p2 = tf.add_paragraph()
    p2.text = line
    p2.font.size = Pt(sz)
    p2.font.color.rgb = clr
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# Bottom summary
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
             "Q-table: 41 rows (states) x 3 columns (actions).  "
             "Each cell = expected future reward for taking that action in that state.",
             font_size=18, color=DIM_WHITE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: What We Expected vs What We Found (Ablation)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "What We Expected vs. What We Found",
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "We tested every combination of dropout-handling strategies:",
             font_size=22, color=LIGHT_GRAY)

# Table
table_data = [
    ("Strategy", "Reward", ""),
    ("Baseline (no dropout)", "-31.80", "reference"),
    ("Naive (always update)", "-45.94", "neutral"),
    ("Confidence alpha only", "-45.66", "neutral"),
    ("State estimator only", "-68.10", "bad"),
    ("Estimator + confidence", "-68.02", "bad"),
    ("Skip when stale", "-43.24", "good"),
]

tbl_shape = slide.shapes.add_table(
    len(table_data), 2, Inches(1.5), Inches(2.0), Inches(7), Inches(4.0)
)
table = tbl_shape.table
table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(2.5)

for row_idx, (strat, reward, style) in enumerate(table_data):
    cell0 = table.cell(row_idx, 0)
    cell1 = table.cell(row_idx, 1)
    cell0.text = strat
    cell1.text = reward

    for cell in [cell0, cell1]:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = "Calibri"
            if row_idx == 0:
                paragraph.font.size = Pt(20)
                paragraph.font.bold = True
                paragraph.font.color.rgb = ACCENT_YELLOW
            else:
                paragraph.font.size = Pt(20)
                paragraph.font.bold = (style == "good")
                if style == "bad":
                    paragraph.font.color.rgb = ACCENT_RED
                elif style == "good":
                    paragraph.font.color.rgb = ACCENT_GREEN
                elif style == "reference":
                    paragraph.font.color.rgb = ACCENT_BLUE
                else:
                    paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER

        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        elif row_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)

# Right side insight
add_rounded_rect(slide, Inches(9), Inches(2.2), Inches(3.8), Inches(3.5),
                 RGBColor(0x8B, 0x30, 0x30))

_, tf = add_text_box(slide, Inches(9.2), Inches(2.4), Inches(3.4), Inches(3.2),
                     "The fancy solution\nmade it WORSE",
                     font_size=26, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "State estimator: -68.10", font_size=18, color=ACCENT_RED,
              alignment=PP_ALIGN.CENTER, space_before=16)
add_paragraph(tf, "vs. doing nothing: -45.94", font_size=18, color=LIGHT_GRAY,
              alignment=PP_ALIGN.CENTER, space_before=8)

# ============================================================
# SLIDE 8: The Winning Strategy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(2),
             '"When you don\'t know,\n don\'t learn."',
             font_size=56, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(2.5),
             "~28% of steps are skipped (83 out of 300 per episode)\n\n"
             "But the remaining 72% write accurate values\n"
             "to the correct Q-table entries.\n\n"
             "A smaller set of clean updates beats\na larger set of noisy ones.",
             font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
             line_spacing=1.4)

# ============================================================
# SLIDE 9: Training Curves (Plot 1)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(7), Inches(0.8),
             "Training: 500 Episodes",
             font_size=40, color=WHITE, bold=True)

plot_path = os.path.join(PLOTS_DIR, "1_training_curves.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(0.3), Inches(1.2),
                             width=Inches(8.5))

# Right side stats
add_rounded_rect(slide, Inches(9.2), Inches(1.5), Inches(3.8), Inches(5),
                 RGBColor(0x2A, 0x2A, 0x45))

_, tf = add_text_box(slide, Inches(9.4), Inches(1.7), Inches(3.4), Inches(4.6),
                     "Last-20 Episode Average",
                     font_size=18, color=ACCENT_YELLOW, bold=True,
                     alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=8, color=WHITE)
add_paragraph(tf, "Baseline", font_size=20, color=ACCENT_GREEN,
              alignment=PP_ALIGN.CENTER, space_before=16, bold=True)
add_paragraph(tf, "-31.80", font_size=28, color=ACCENT_GREEN,
              alignment=PP_ALIGN.CENTER, bold=True)
add_paragraph(tf, "", font_size=8, color=WHITE)
add_paragraph(tf, "Dropout-Aware", font_size=20, color=ACCENT_BLUE,
              alignment=PP_ALIGN.CENTER, space_before=16, bold=True)
add_paragraph(tf, "-43.24", font_size=28, color=ACCENT_BLUE,
              alignment=PP_ALIGN.CENTER, bold=True)
add_paragraph(tf, "", font_size=8, color=WHITE)
add_paragraph(tf, "Naive", font_size=20, color=ACCENT_RED,
              alignment=PP_ALIGN.CENTER, space_before=16, bold=True)
add_paragraph(tf, "-45.94", font_size=28, color=ACCENT_RED,
              alignment=PP_ALIGN.CENTER, bold=True)
add_paragraph(tf, "", font_size=8, color=WHITE)
add_paragraph(tf, "~6% improvement from\nskipping stale updates",
              font_size=16, color=LIGHT_GRAY,
              alignment=PP_ALIGN.CENTER, space_before=20)

# ============================================================
# SLIDE 10: Trajectories (Plot 4)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "What Does It Look Like?",
             font_size=40, color=WHITE, bold=True)

plot_path = os.path.join(PLOTS_DIR, "4_trajectories.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(2.5), Inches(1.0),
                             width=Inches(8.2))

add_text_box(slide, Inches(0.5), Inches(1.2), Inches(2.2), Inches(5),
             "Top:\nBaseline hugs\nthe target\n\n"
             "Middle:\nDropout agents\nshow overshoots\nat peaks\n\n"
             "Bottom:\nRed = packet\nwas lost",
             font_size=16, color=LIGHT_GRAY, line_spacing=1.3)

# ============================================================
# SLIDE 11: The Second Surprise — Evaluation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "The Second Surprise",
             font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "At evaluation time (greedy policy, no learning):",
             font_size=24, color=LIGHT_GRAY)

# Eval results
tbl_shape = slide.shapes.add_table(
    4, 3, Inches(2), Inches(2.2), Inches(9), Inches(2.4)
)
table = tbl_shape.table
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(3)
table.columns[2].width = Inches(3)

eval_data = [
    ("Agent", "Mean Tracking Error", "Total Reward"),
    ("Baseline", "0.0858", "-25.75"),
    ("Naive", "0.1236", "-37.09"),
    ("Dropout-Aware", "0.1236", "-37.09"),
]

for row_idx, row_data in enumerate(eval_data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_YELLOW
            elif row_idx >= 2 and col_idx >= 1:
                p.font.color.rgb = ACCENT_ORANGE
                p.font.bold = True
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        elif row_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)

# Highlight: identical!
add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.6),
             "Naive and Dropout-Aware: IDENTICAL evaluation scores",
             font_size=28, color=ACCENT_ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)

_, tf = add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(1.5),
                     "With 41 states and 3 actions, both Q-tables converge "
                     "to the same argmax.",
                     font_size=20, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "The dropout-aware advantage is learning efficiency, "
              "not the final policy.",
              font_size=20, color=WHITE, bold=True,
              alignment=PP_ALIGN.CENTER, space_before=8)

# ============================================================
# SLIDE 12: Transition to IRL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.5),
             "But wait...",
             font_size=44, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.5),
             "We hand-coded the reward: -|distance|",
             font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(1.5),
             "How did we know that was right?",
             font_size=36, color=ACCENT_RED, bold=True,
             alignment=PP_ALIGN.CENTER)

_, tf = add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(1.5),
                     "What if we could discover the reward automatically—",
                     font_size=24, color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "just by watching an expert do the task?",
              font_size=24, color=LIGHT_GRAY,
              alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: IRL Concept
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Inverse Reinforcement Learning",
             font_size=40, color=WHITE, bold=True)

# Two arrows showing the flip
# RL arrow
add_rounded_rect(slide, Inches(1), Inches(1.8), Inches(5.2), Inches(2.2),
                 RGBColor(0x2A, 0x2A, 0x45))
_, tf = add_text_box(slide, Inches(1.2), Inches(1.9), Inches(4.8), Inches(2),
                     "Regular RL",
                     font_size=24, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "Given: reward function", font_size=20, color=WHITE,
              space_before=8)
add_paragraph(tf, "Find: optimal policy", font_size=20, color=ACCENT_GREEN,
              space_before=8, bold=True)

# IRL arrow
add_rounded_rect(slide, Inches(7), Inches(1.8), Inches(5.2), Inches(2.2),
                 RGBColor(0x2A, 0x2A, 0x45))
_, tf = add_text_box(slide, Inches(7.2), Inches(1.9), Inches(4.8), Inches(2),
                     "Inverse RL",
                     font_size=24, color=ACCENT_ORANGE, bold=True)
add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "Given: expert demonstrations", font_size=20, color=WHITE,
              space_before=8)
add_paragraph(tf, "Find: reward function", font_size=20, color=ACCENT_ORANGE,
              space_before=8, bold=True)

# Analogy
add_rounded_rect(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(2),
                 RGBColor(0x33, 0x33, 0x55))

_, tf = add_text_box(slide, Inches(1.8), Inches(5.0), Inches(9.5), Inches(1.8),
                     "The Analogy",
                     font_size=22, color=ACCENT_YELLOW, bold=True)
add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "You watch a top student ace every exam.",
              font_size=20, color=WHITE, space_before=8)
add_paragraph(tf, "You don't copy their answers.",
              font_size=20, color=WHITE, space_before=4)
add_paragraph(tf, "You figure out what textbook they're studying from.",
              font_size=20, color=ACCENT_ORANGE, bold=True, space_before=4)

# ============================================================
# SLIDE 14: The 5 Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "5 Hypotheses About What the Expert Cares About",
             font_size=36, color=WHITE, bold=True)

features_data = [
    ("f0: Bullseye", "1.0 if within 0.1 units of target",
     "Does the expert care about precision?", ACCENT_GREEN),
    ("f1: Linear", "-|distance|",
     "Is the penalty proportional to distance?", ACCENT_BLUE),
    ("f2: Quadratic", "-distance²",
     "Is the penalty proportional to distance²?", ACCENT_RED),
    ("f3: Danger", "-1.0 if beyond 75% of range",
     "Does the expert avoid edges specifically?", ACCENT_ORANGE),
    ("f4: Bias", "constant 1.0",
     "Is there a constant cost per step?", DIM_WHITE),
]

for i, (name, formula, question, color) in enumerate(features_data):
    y = Inches(1.4 + i * 1.15)
    add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.95),
                     RGBColor(0x2A, 0x2A, 0x45))
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(2.5), Inches(0.7),
                 name, font_size=22, color=color, bold=True)
    add_text_box(slide, Inches(3.8), y + Inches(0.1), Inches(3.5), Inches(0.7),
                 formula, font_size=20, color=WHITE, font_name="Courier New")
    add_text_box(slide, Inches(7.5), y + Inches(0.1), Inches(4.5), Inches(0.7),
                 question, font_size=18, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
             "We give IRL these building blocks. It discovers which ones matter.",
             font_size=20, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15: IRL Loop
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "The IRL Loop: 200 Iterations",
             font_size=40, color=WHITE, bold=True)

# IRL diagram on the right
if os.path.exists(IRL_DIAGRAM_PATH):
    slide.shapes.add_picture(IRL_DIAGRAM_PATH, Inches(7.3), Inches(1.2),
                             width=Inches(5.5))

# Loop steps on the left
loop_steps = [
    ("1", "Train Q-learner (100 episodes)\nreward = weights · features",
     ACCENT_BLUE),
    ("2", "Run policy (10 episodes)\nMeasure average feature vector",
     ACCENT_GREEN),
    ("3", "gradient = expert - learner\n\"Where is the learner falling short?\"",
     ACCENT_ORANGE),
    ("4", "weights += 0.01 × gradient\n\"Push reward toward expert behavior\"",
     ACCENT_YELLOW),
]

for i, (num, desc, color) in enumerate(loop_steps):
    y = Inches(1.4 + i * 1.3)
    # Number circle
    add_rounded_rect(slide, Inches(0.8), y, Inches(0.6), Inches(0.6),
                     color, text=num, font_size=24, font_color=BG_DARK)
    # Description
    add_text_box(slide, Inches(1.6), y + Inches(0.0), Inches(5.5), Inches(1.1),
                 desc, font_size=20, color=WHITE, line_spacing=1.3)

# Bottom: convergence insight
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.7),
             "Repeat until the learner's behavior matches the expert's "
             "→ reward function found",
             font_size=20, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 16: IRL Convergence (Plot 5)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(7), Inches(0.8),
             "What IRL Discovered",
             font_size=40, color=WHITE, bold=True)

plot_path = os.path.join(PLOTS_DIR, "5_irl_convergence.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(0.3), Inches(1.2),
                             width=Inches(8.5))

# Right side: final weights
add_rounded_rect(slide, Inches(9.2), Inches(1.5), Inches(3.8), Inches(5.2),
                 RGBColor(0x2A, 0x2A, 0x45))

_, tf = add_text_box(slide, Inches(9.4), Inches(1.7), Inches(3.4), Inches(5),
                     "Final Learned Weights",
                     font_size=18, color=ACCENT_YELLOW, bold=True,
                     alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=8, color=WHITE)

weight_items = [
    ("Quadratic", "10.000", ACCENT_RED, True),
    ("Bullseye", "0.806", ACCENT_GREEN, False),
    ("Linear", "0.378", ACCENT_BLUE, False),
    ("Danger", "0.010", ACCENT_ORANGE, False),
    ("Bias", "0.000", DIM_WHITE, False),
]

for name, val, clr, dominant in weight_items:
    add_paragraph(tf, name, font_size=18, color=clr,
                  alignment=PP_ALIGN.CENTER, space_before=14, bold=True)
    add_paragraph(tf, val, font_size=24 if dominant else 20,
                  color=clr, alignment=PP_ALIGN.CENTER, bold=dominant)

add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "Quadratic dominates—", font_size=16, color=LIGHT_GRAY,
              alignment=PP_ALIGN.CENTER, space_before=12)
add_paragraph(tf, "being far is very costly", font_size=16, color=LIGHT_GRAY,
              alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 17: 3x2 Comparison (Plot 6)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "The Proof: Hand-Coded vs. IRL-Learned",
             font_size=40, color=WHITE, bold=True)

# 3x2 table
tbl_shape = slide.shapes.add_table(
    4, 3, Inches(1), Inches(1.5), Inches(7), Inches(3.0)
)
table = tbl_shape.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.25)
table.columns[2].width = Inches(2.25)

compare_data = [
    ("Agent", "Hand-coded Error", "IRL-learned Error"),
    ("Baseline", "0.0858", "0.0858"),
    ("Naive", "0.1236", "0.1236"),
    ("Dropout-Aware", "0.1236", "0.1236"),
]

for row_idx, row_data in enumerate(compare_data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_YELLOW
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        elif row_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)

add_text_box(slide, Inches(1), Inches(4.7), Inches(7), Inches(0.8),
             "Every cell matches. IRL recovered the reward perfectly.",
             font_size=24, color=ACCENT_GREEN, bold=True)

# Plot on the right
plot_path = os.path.join(PLOTS_DIR, "6_irl_comparison.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(8.2), Inches(1.5),
                             width=Inches(4.8))

add_text_box(slide, Inches(8.2), Inches(5.5), Inches(4.8), Inches(0.8),
             "Green (hand-coded) and orange (IRL)\noverlap exactly",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 2D EXTENSION SLIDES
# ============================================================
PLOTS_2D_DIR = os.path.join(OUTPUT_DIR, "experiment_2d", "output_plots")

# --- SLIDE: Transition — "But was it too easy?" ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.5),
             "But was our problem too easy?",
             font_size=44, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)

_, tf = add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(3.5),
                     "In 1D, Naive and Dropout-Aware converged to the same policy.",
                     font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "41 states, 3 actions — too simple to corrupt.",
              font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
              space_before=12)
add_paragraph(tf, "", font_size=10, color=WHITE)
add_paragraph(tf, "What happens when we scale up?",
              font_size=30, color=ACCENT_RED, bold=True,
              alignment=PP_ALIGN.CENTER, space_before=16)

# --- SLIDE: 2D Environment Setup ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Scaling Up: 2D Figure-Eight Tracking",
             font_size=40, color=WHITE, bold=True)

# Left: 1D vs 2D comparison
add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2),
                 RGBColor(0x2A, 0x2A, 0x45))

_, tf = add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5.1), Inches(5),
                     "", font_size=20, color=ACCENT_YELLOW, bold=True)

tbl_2d = slide.shapes.add_table(
    6, 3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.0)
)
table = tbl_2d.table
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(1.8)
table.columns[2].width = Inches(1.9)

setup_data = [
    ("", "1D (before)", "2D (now)"),
    ("States", "41", "441 (21x21)"),
    ("Actions", "3", "9 (8 dirs + stay)"),
    ("Q-table", "123 entries", "3,969 entries"),
    ("Target", "Sine wave", "Figure-eight"),
    ("Training", "500 episodes", "800 episodes"),
]

for row_idx, row_data in enumerate(setup_data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_YELLOW
            elif col_idx == 2:
                p.font.color.rgb = ACCENT_GREEN
                p.font.bold = True
            elif col_idx == 1:
                p.font.color.rgb = LIGHT_GRAY
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        elif row_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)

# Right: why this matters
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                 RGBColor(0x33, 0x33, 0x55))

_, tf = add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.4), Inches(5),
                     "Why 2D changes everything",
                     font_size=22, color=ACCENT_YELLOW, bold=True)
add_paragraph(tf, "", font_size=8, color=WHITE)
add_paragraph(tf, "32x larger Q-table",
              font_size=20, color=ACCENT_GREEN, bold=True, space_before=14)
add_paragraph(tf, "Many states visited rarely — corrupted\nupdates can flip the best action",
              font_size=16, color=LIGHT_GRAY, space_before=4)
add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "3x more actions",
              font_size=20, color=ACCENT_GREEN, bold=True, space_before=14)
add_paragraph(tf, "Choosing between 8 directions is\nambiguous — noise promotes wrong ones",
              font_size=16, color=LIGHT_GRAY, space_before=4)
add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "2D stale states",
              font_size=20, color=ACCENT_GREEN, bold=True, space_before=14)
add_paragraph(tf, "Wrong in both X and Y simultaneously\n— bigger mistakes, harder to recover",
              font_size=16, color=LIGHT_GRAY, space_before=4)

# --- SLIDE: 2D Trajectories ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
             "2D Trajectories: The Difference Is Visible",
             font_size=38, color=WHITE, bold=True)

plot_path = os.path.join(PLOTS_2D_DIR, "2_trajectories_2d.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.1),
                             width=Inches(12.3))

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(1),
             "Dropout-Aware (blue) traces the figure-eight more smoothly. "
             "Naive (red) shows sharper deviations from corrupted Q-table entries.",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# --- SLIDE: 2D Eval Results ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.8),
             "2D Results: The Gap Is Real",
             font_size=40, color=WHITE, bold=True)

# Results table
tbl_shape = slide.shapes.add_table(
    4, 3, Inches(0.8), Inches(1.5), Inches(7), Inches(3.0)
)
table = tbl_shape.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.25)
table.columns[2].width = Inches(2.25)

eval_2d_data = [
    ("Agent", "Mean Track Error", "Total Reward"),
    ("Baseline (no drop)", "0.2373", "-85.44"),
    ("Naive (30% drop)", "0.4413", "-158.87"),
    ("Dropout-Aware (30%)", "0.3730", "-134.29"),
]

for row_idx, row_data in enumerate(eval_2d_data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(22)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_YELLOW
            elif row_idx == 3:
                p.font.color.rgb = ACCENT_GREEN
                p.font.bold = True
            elif row_idx == 1:
                p.font.color.rgb = ACCENT_ORANGE
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        elif row_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)

# Highlight box
add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(7), Inches(1.8),
                 RGBColor(0x1E, 0x3A, 0x5C))

_, tf = add_text_box(slide, Inches(1.0), Inches(4.9), Inches(6.6), Inches(1.6),
                     "Dropout-Aware beats Naive by 15.5%",
                     font_size=26, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "", font_size=6, color=WHITE)
add_paragraph(tf, "Unlike 1D, the policies are now DIFFERENT.",
              font_size=20, color=WHITE, space_before=8)
add_paragraph(tf, "Corrupted Q-table entries produce wrong actions\n"
              "in rarely-visited states.",
              font_size=18, color=LIGHT_GRAY, space_before=6)

# Bar chart on the right
plot_path = os.path.join(PLOTS_2D_DIR, "4_eval_comparison.png")
if os.path.exists(plot_path):
    slide.shapes.add_picture(plot_path, Inches(8.0), Inches(1.3),
                             width=Inches(5))

# --- SLIDE: 1D vs 2D side-by-side ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
             "1D vs 2D: Complexity Reveals the Truth",
             font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# 1D box
add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2),
                 RGBColor(0x2A, 0x2A, 0x45))

_, tf = add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5.2), Inches(0.5),
                     "1D Sine Wave (Simple)",
                     font_size=24, color=ACCENT_ORANGE, bold=True,
                     alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5.2), Inches(1),
             "123 Q-table entries",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tbl_1d = slide.shapes.add_table(
    4, 2, Inches(1.2), Inches(3.0), Inches(4.8), Inches(2.5)
)
table = tbl_1d.table
table.columns[0].width = Inches(2.4)
table.columns[1].width = Inches(2.4)

data_1d = [
    ("Agent", "Eval Error"),
    ("Baseline", "0.0858"),
    ("Naive", "0.1236"),
    ("Dropout-Aware", "0.1236"),
]

for row_idx, row_data in enumerate(data_1d):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_YELLOW
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)

add_text_box(slide, Inches(1.0), Inches(5.7), Inches(5.2), Inches(0.8),
             "Naive = Dropout-Aware\nSame policy, no visible difference",
             font_size=18, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER,
             bold=True)

# 2D box
add_rounded_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                 RGBColor(0x2A, 0x2A, 0x45))

_, tf = add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.5),
                     "2D Figure-Eight (Complex)",
                     font_size=24, color=ACCENT_GREEN, bold=True,
                     alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(1),
             "3,969 Q-table entries (32x larger)",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tbl_2d = slide.shapes.add_table(
    4, 2, Inches(7.2), Inches(3.0), Inches(5.0), Inches(2.5)
)
table = tbl_2d.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.5)

data_2d = [
    ("Agent", "Eval Error"),
    ("Baseline (no drop)", "0.2373"),
    ("Naive (30% drop)", "0.4413"),
    ("Dropout-Aware", "0.3730"),
]

for row_idx, row_data in enumerate(data_2d):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = ACCENT_YELLOW
            elif row_idx == 3:
                p.font.color.rgb = ACCENT_GREEN
                p.font.bold = True
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        if row_idx == 0:
            cell.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x55)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x42)

add_text_box(slide, Inches(7.0), Inches(5.7), Inches(5.4), Inches(0.8),
             "Dropout-Aware BEATS Naive by 15.5%\nComplexity reveals the real advantage",
             font_size=18, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER,
             bold=True)

# ============================================================
# SLIDE: Three Lessons (updated number)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
             "Three Takeaways",
             font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

lessons = [
    ("1", "When in doubt, don't act.",
     "Acting on bad information is worse than waiting.\n"
     "The dropout-aware agent improved by simply skipping uncertain updates.",
     ACCENT_GREEN),
    ("2", "Simple beats clever.",
     "A sophisticated state estimator made things worse.\n"
     "The best strategy was the most obvious one.",
     ACCENT_BLUE),
    ("3", "You don't need to define the goal.",
     "Inverse RL discovered the reward from expert demonstrations alone.\n"
     "The objective reveals itself through behavior.",
     ACCENT_ORANGE),
]

for i, (num, title, desc, color) in enumerate(lessons):
    y = Inches(1.5 + i * 1.9)
    add_rounded_rect(slide, Inches(1), y, Inches(0.8), Inches(0.8),
                     color, text=num, font_size=32, font_color=BG_DARK)
    add_text_box(slide, Inches(2.1), y + Inches(0.0), Inches(10), Inches(0.6),
                 title, font_size=28, color=color, bold=True)
    add_text_box(slide, Inches(2.1), y + Inches(0.7), Inches(10), Inches(1.0),
                 desc, font_size=18, color=LIGHT_GRAY, line_spacing=1.3)

# ============================================================
# SLIDE 19: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(2),
             "The next time a robot loses signal\n"
             "mid-operation...",
             font_size=36, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(2.5),
             "the smartest thing it can do\n"
             "might be the simplest:",
             font_size=36, color=WHITE,
             alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(1.5),
             "Put down the scalpel.\nWait for a clear picture.\nThen cut.",
             font_size=40, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER, line_spacing=1.4)

# ============================================================
# SLIDE 20: Thank You / Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
             "Thank You",
             font_size=56, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             "Questions?",
             font_size=36, color=ACCENT_BLUE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             "Doron Peleg  |  doron.peleg@gmail.com",
             font_size=20, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Code: github.com/dpartuk/RL",
             font_size=18, color=DIM_WHITE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Based on: Fan et al. (2025) \"Inverse RL for Discrete-Time Systems "
             "With Data Dropouts\" — IEEE Trans. Cybernetics, Vol. 55, No. 4",
             font_size=14, color=DIM_WHITE,
             alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
prs.save(OUT_PATH)
print(f"Presentation saved to: {OUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
